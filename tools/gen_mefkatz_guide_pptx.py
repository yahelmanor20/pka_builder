# -*- coding: utf-8 -*-
"""
מחולל מצגת "איך לכתוב לו"ז נכון ביומן — מדריך למפקץ" (PPTX עברית RTL, מעוצב).
אותו תוכן כמו דף הכללים (gen_mefkatz_guide.py), בפריסת שקופיות מקצועית:
פס כותרת צבעוני + תג מספר בכל שקופית · כרטיסים · תיבות נכון/שגוי · כותרת תחתונה.
כל כלל נגזר מהפרסר ב-index.html. דורש python-pptx.
פלט: מדריך_כתיבת_לוז_ביומן.pptx בשורש הפרויקט.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(__file__)
OUT = os.path.join(HERE, "..", "מדריך_כתיבת_לוז_ביומן.pptx")

FONT = "David"
BLUE = RGBColor(0x1F, 0x4E, 0x79)
BLUE2 = RGBColor(0x2E, 0x6D, 0xA4)
DARK = RGBColor(0x2A, 0x2A, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE9, 0xF0, 0xF7)
GRAYTX = RGBColor(0x6B, 0x6B, 0x6B)
GREEN = RGBColor(0x2E, 0x7D, 0x32); GREEN_BG = RGBColor(0xE6, 0xF3, 0xE6)
RED = RGBColor(0xB0, 0x00, 0x00); RED_BG = RGBColor(0xFB, 0xE6, 0xE1)
GOLD = RGBColor(0x8A, 0x66, 0x00); GOLD_BG = RGBColor(0xFF, 0xF3, 0xCC)
CARD_BG = RGBColor(0xF4, 0xF6, 0xF9)

SW, SH = Inches(13.333), Inches(7.5)

COLORS = [
    (RGBColor(0xD5, 0x00, 0x00), WHITE, "עגבנייה (Tomato)",   "פלוגה א׳"),
    (RGBColor(0xF4, 0x51, 0x1E), WHITE, "מנדרינה (Tangerine)", "פלוגה ב׳"),
    (RGBColor(0xF6, 0xBF, 0x26), DARK,  "בננה (Banana)",       "פלוגה ג׳"),
    (RGBColor(0x0B, 0x80, 0x43), WHITE, "בזיליקום (Basil)",     "פלוגה ד׳"),
    (RGBColor(0xE6, 0x7C, 0x73), DARK,  "פלמינגו (Flamingo)",   "לו״ז מקביל → נספח"),
    (RGBColor(0x61, 0x61, 0x61), WHITE, "גרפיט (Graphite)",     "פלס״ם → נספח פלס״ם"),
]


# ---------- כלי עיצוב ----------
def _rtl(p):
    p._p.get_or_add_pPr().set("rtl", "1")


def _set(run, size, bold=False, color=DARK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rPr.append(rPr.makeelement(qn("a:cs"), {"typeface": FONT}))


def _noshadow(shape):
    el = shape._element.spPr
    el.append(el.makeelement(qn("a:effectLst"), {}))


def rect(slide, l, t, w, h, fill, line=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    _noshadow(shp)
    if rounded:
        try:
            shp.adjustments[0] = 0.07
        except Exception:
            pass
    return shp


def text(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.RIGHT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        txt, size, bold, color = ln[0], ln[1], ln[2], ln[3]
        al = ln[4] if len(ln) > 4 else align
        sp = ln[5] if len(ln) > 5 else 4
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al; p.space_after = Pt(sp); _rtl(p)
        r = p.add_run(); r.text = txt; _set(r, size, bold, color)
    return box


def fill_text(shape, lines, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    for i, ln in enumerate(lines):
        txt, size, bold, color = ln[0], ln[1], ln[2], ln[3]
        al = ln[4] if len(ln) > 4 else PP_ALIGN.RIGHT
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al; p.space_after = Pt(3); _rtl(p)
        r = p.add_run(); r.text = txt; _set(r, size, bold, color)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, num, head):
    """פס כותרת עליון + תג מספר עגול."""
    rect(slide, 0, 0, 13.333, 1.15, BLUE)
    rect(slide, 0, 1.15, 13.333, 0.07, BLUE2)
    # תג מספר בצד ימין (RTL)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.0), Inches(0.27), Inches(0.62), Inches(0.62))
    badge.fill.solid(); badge.fill.fore_color.rgb = WHITE; badge.line.fill.background(); _noshadow(badge)
    fill_text(badge, [(str(num), 24, True, BLUE, PP_ALIGN.CENTER)])
    text(slide, 0.5, 0.18, 11.2, 0.85, [(head, 27, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)


def footer(slide):
    text(slide, 0.4, 7.04, 12.5, 0.4,
         [("מדריך למפקץ · איך לכתוב לו״ז נכון ביומן · בא״פ 8222", 10, False, GRAYTX, PP_ALIGN.LEFT)])


def good_bad(slide, top, good, bad):
    w, gap = 5.9, 0.4
    x_bad = 0.7
    x_good = x_bad + w + gap
    for x, label, lab_c, bg, txt in [
        (x_good, "✓  נכון", GREEN, GREEN_BG, good),
        (x_bad, "✗  שגוי", RED, RED_BG, bad),
    ]:
        rect(slide, x, top, w, 1.55, bg, rounded=True)
        text(slide, x + 0.15, top + 0.1, w - 0.3, 0.5, [(label, 18, True, lab_c, PP_ALIGN.CENTER)])
        text(slide, x + 0.2, top + 0.62, w - 0.4, 0.85, [(txt, 16, False, DARK, PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)


def content_slide(prs, num, head, body_lines, gb=None, callout=None, callout_kind="gold"):
    s = blank(prs); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    header(s, num, head)
    text(s, 0.7, 1.55, 11.9, 3.6, body_lines, anchor=MSO_ANCHOR.TOP)
    y = 5.0
    if callout:
        kc = {"gold": (GOLD_BG, GOLD), "red": (RED_BG, RED)}[callout_kind]
        rect(s, 0.7, y, 11.93, 0.85, kc[0], rounded=True)
        text(s, 0.9, y + 0.12, 11.5, 0.6, [(callout, 16, True, kc[1], PP_ALIGN.CENTER)],
             anchor=MSO_ANCHOR.MIDDLE)
        y += 1.1
    if gb:
        good_bad(s, max(y, 5.0), gb[0], gb[1])
    footer(s)
    return s


def bl(txt, bold=False, color=DARK, size=19, sp=7):
    return (txt, size, bold, color, PP_ALIGN.RIGHT, sp)


# ---------- בניית המצגת ----------
def build(path):
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    # ---- פתיחה ----
    s = blank(prs); s.background.fill.solid(); s.background.fill.fore_color.rgb = BLUE
    rect(s, 0, 5.7, 13.333, 1.8, BLUE2)
    text(s, 0.8, 2.2, 11.7, 3.0, [
        ("איך לכתוב לו״ז נכון ביומן", 50, True, WHITE, PP_ALIGN.CENTER, 10),
        ("מדריך למפקץ · בא״פ 8222 · ביה״ס ללחימה", 24, False, LIGHT, PP_ALIGN.CENTER, 6),
    ], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.5, 5.95, 10.3, 1.2, [
        ("הכלי בונה את הפק״א מהיומן אוטומטית.", 19, True, WHITE, PP_ALIGN.CENTER, 3),
        ("כתוב את האירועים לפי 8 הכללים הבאים — והוא יקרא אותם נכון.", 17, False, LIGHT, PP_ALIGN.CENTER, 3),
    ], anchor=MSO_ANCHOR.MIDDLE)

    # ---- 1: צבע = פלוגה (טבלה) ----
    s = blank(prs); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    header(s, 1, "הכי חשוב: הצבע קובע את הפלוגה")
    rect(s, 0.7, 1.45, 11.93, 0.75, RED_BG, rounded=True)
    text(s, 0.9, 1.55, 11.5, 0.55,
         [("הצבע גובר על הכותרת — אירוע בלי צבע לא משויך לאף פלוגה!", 18, True, RED, PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)
    rows = len(COLORS) + 1
    tb = s.shapes.add_table(rows, 2, Inches(1.7), Inches(2.45), Inches(9.9), Inches(4.0)).table
    tb.first_row = False; tb.horz_banding = False
    tb.columns[0].width = Inches(4.3); tb.columns[1].width = Inches(5.6)
    for j, h in enumerate(["צבע האירוע ביומן", "כך הכלי משייך אותו"]):
        c = tb.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = BLUE
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; _rtl(p)
        r = p.add_run(); r.text = h; _set(r, 16, True, WHITE)
    for i, (fl, tc, name, meaning) in enumerate(COLORS, 1):
        c0 = tb.cell(i, 0); c0.fill.solid(); c0.fill.fore_color.rgb = fl
        c0.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c0.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; _rtl(p)
        r = p.add_run(); r.text = name; _set(r, 16, True, tc)
        c1 = tb.cell(i, 1); c1.fill.solid(); c1.fill.fore_color.rgb = WHITE
        c1.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c1.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; _rtl(p)
        r = p.add_run(); r.text = "  " + meaning; _set(r, 15, False, DARK)
    text(s, 0.7, 6.6, 11.9, 0.45,
         [("ביומן: פותחים את האירוע ← לוחצים על עיגול הצבע ← בוחרים את הצבע הנכון.", 14, False, GRAYTX)])
    footer(s)

    # ---- 2: כותרת ----
    content_slide(prs, 2, "כותרת אירוע פלוגתי", [
        bl("הפורמט:   פלוגה א׳ - שם הפעילות", True, BLUE, 22),
        bl("הטקסט שאחרי המקף = שם הפעילות שיופיע בלו״ז.", False, DARK, 19),
        bl("הכלי סובלני (פל׳ / פלוגה / מקף / נקודתיים) — אבל כתוב מסודר.", False, GRAYTX, 17),
    ], gb=("פלוגה א׳ - מטווח קלעים", "מטווח קלעים — בלי פלוגה ובלי צבע"))

    # ---- 3: שעה ----
    content_slide(prs, 3, "שעה — בשדה השעה, לא בכותרת", [
        bl("מלא שעת התחלה וסיום אמיתיות באירוע.", True, DARK, 20),
        bl("מהן הכלי בונה את עמודת השעה בלו״ז.", False, DARK, 19),
        bl("אירוע ״כל היום״ או בלי שעה → אזהרה ״ללא שעה״.", False, RED, 18),
    ], gb=("אירוע מתוזמן 08:00–11:00", "אירוע ״כל היום״, או שעה בתוך הכותרת"))

    # ---- 4: מיקום ----
    content_slide(prs, 4, "מיקום", [
        bl("מלא את שדה המיקום של האירוע ביומן.", True, DARK, 20),
        bl("שדה ריק → ״מיקום לא ידוע״ + אזהרה צהובה.", False, RED, 18),
        bl("אפשר גם בתיאור, בשורה:   מיקום: <שם המקום>", False, DARK, 18),
    ], gb=("שדה מיקום: ״שטח 9א׳״", "שדה מיקום ריק; המקום רק בכותרת"))

    # ---- 5: תיאור ----
    s = blank(prs); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    header(s, 5, "תיאור — שדות מסודרים")
    text(s, 0.7, 1.45, 11.9, 0.6,
         [("כל שדה בשורה נפרדת, עם נקודתיים. הכלי מזהה בדיוק את שלוש המילים:", 19, True, DARK)])
    for i, (k, v) in enumerate([
        ("גורם מעביר:", "שם המעביר — חסר → אזהרה ״חסר גורם מעביר״"),
        ("הערות:", "טקסט חופשי שיופיע בעמודת ההערות"),
        ("מיקום:", "שם מקום — דורס את שדה המיקום"),
    ]):
        y = 2.25 + i * 1.0
        rect(s, 2.4, y, 9.0, 0.82, CARD_BG, rounded=True)
        rect(s, 10.7, y, 0.7, 0.82, BLUE, rounded=True)
        text(s, 8.0, y + 0.13, 2.6, 0.55, [(k, 18, True, BLUE, PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 2.6, y + 0.13, 5.2, 0.55, [(v, 15, False, DARK)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.7, 5.55, 11.9, 1.0, [
        ("כל שורה שאינה בפורמט הזה — נדבקת אוטומטית לשדה ההערות.", 16, False, GRAYTX, PP_ALIGN.RIGHT, 5),
        ("באירוע טכניקות/תרגול — פירוט בתיאור נדבק בסוגריים לשם הפעילות.", 16, False, GRAYTX, PP_ALIGN.RIGHT, 5),
    ])
    footer(s)

    # ---- 6: מקבילות ----
    content_slide(prs, 6, "פעילויות מקבילות (לא בלו״ז היומי)", [
        bl("חובשים · תאג״ד · סמב״צים · הסמכות · רענונים — צבע פלמינגו → נספח מקביל.", True, DARK, 19),
        bl("פלס״ם — צבע גרפיט → נספח פלס״ם נפרד.", True, DARK, 19),
    ], callout="חריג: ״מרחב חכם״ תמיד נשאר בלו״ז היומי — גם אם צבעת אותו פלמינגו.",
        callout_kind="gold")

    # ---- 7: יום מפוצל ----
    content_slide(prs, 7, "יום מפוצל (2+ פלוגות באותו יום)", [
        bl("מזהה 2 פלוגות שונות או יותר ביום → הטבלה הופכת אוטומטית לעמודות פלוגה.", True, DARK, 19),
        bl("", False, DARK, 8),
        bl("פעילות לכל הגדוד ביום מפוצל:", True, BLUE, 20),
        bl("השאר אותה בלי צבע פלוגה — כך תופיע כשורה רוחבית אחת לכל הגדוד.", False, DARK, 19),
    ])

    # ---- 8: צהוב + אל תעשה ----
    s = blank(prs); s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    header(s, 8, "לפני שסיימת — בדוק את הצהוב")
    rect(s, 0.7, 1.55, 11.93, 1.35, GOLD_BG, rounded=True)
    text(s, 0.9, 1.7, 11.5, 1.1, [
        ("כל אזהרה צהובה בתצוגה = משהו לתקן ביומן:", 20, True, GOLD, PP_ALIGN.CENTER, 6),
        ("ללא שעה   ·   מיקום לא ידוע   ·   חסר גורם מעביר", 20, True, DARK, PP_ALIGN.CENTER, 4),
    ], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 0.7, 3.5, 11.93, 2.7, RED_BG, rounded=True)
    text(s, 0.9, 3.7, 11.5, 2.4, [
        ("✗  אל תעשה", 26, True, RED, PP_ALIGN.CENTER, 10),
        ("שעה בכותרת במקום בשדה", 19, False, DARK, PP_ALIGN.CENTER, 5),
        ("אירוע בלי צבע   ·   מיקום ריק", 19, False, DARK, PP_ALIGN.CENTER, 5),
        ("שיוך פלוגה לפי כותרת בלבד   ·   פעילות מקבילה בצבע פלוגה", 19, False, DARK, PP_ALIGN.CENTER, 5),
    ], anchor=MSO_ANCHOR.MIDDLE)
    footer(s)

    prs.save(path)
    print("נשמר:", os.path.abspath(path), "—", len(prs.slides._sldIdLst), "שקופיות")


if __name__ == "__main__":
    build(OUT)
